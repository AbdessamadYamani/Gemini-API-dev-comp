
from webbrowser import Chrome
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain.tools import tool
from langchain_community.vectorstores import chroma
from langchain_openai import OpenAIEmbeddings
import os
import shutil
import json
import openpyxl  # For XLSX
import pandas as pd
from docx import Document  # For DOCX
import csv
import shutil
from pptx import Presentation
import gspread
from docx import Document  # For DOCX
from langchain_community.vectorstores import FAISS
from langchain_google_genai import GoogleGenerativeAIEmbeddings
import warnings
import markdown
# from services import llama_parse_sce
import sys
# Ignorer les avertissements de type FutureWarning
warnings.filterwarnings('ignore', category=FutureWarning)
sys.path.append(os.path.abspath(os.path.join(os.path.dirname(__file__), '..')))
import settings.api_configs as config


oldest_transcript = ""
class DocClass:
    
    @tool("Doc Loader From All Types")
    def load_content_from_all_types(query: str):
        """Charge et traite le contenu des documents"""
        global_retriever_results=[]
        textual_content_from_all_resources_list = []
        textual_content_from_pdf_list=[]
        textual_content_from_md_list=[]
        textual_content_from_json_list=[]
        textual_content_from_xlsx_list=[]
        textual_content_from_word_list=[]
        textual_content_from_csv_list=[]
        textual_content_from_pptx_list=[]
        textual_content_from_txt_list=[]
        directory_path = config.directory_RAG_path
        # parsed_pdf_files=config.parsed_pdf_files
        # llama_parse_sce.parse_pdf_to_md(directory_path,directory_path,parsed_pdf_files)
        for filename in os.listdir(directory_path):
            file_path = os.path.join(directory_path, filename)
            
#process txt files 
            if filename.endswith('.txt'):
                with open(file_path, 'r', encoding='utf-8') as file:
                    textual_content_from_txt_list.append(file.read())
                    print(f"txt ::Name: {os.path.basename(file.name)}")
#process pdf files 
            # elif filename.endswith('.pdf'):
            #     with open(file_path, 'rb') as file:
            #         pdf_reader = PyPDF2.PdfReader(file)
            #         num_pages = len(pdf_reader.pages)
            #         tables = tabula.read_pdf(file, pages='all', multiple_tables=True)
            #         cumulated_text_from_pdf = ""

            #         for page_num in range(num_pages):
            #             page = pdf_reader.pages[page_num]
            #             cumulated_text_from_pdf += page.extract_text()
            #         textual_content_from_pdf_list.append(cumulated_text_from_pdf)

            #         cumulated_tables_as_text_from_pdf = ""
            #         for table in tables:
            #             if isinstance(table, pd.DataFrame):
            #                 cumulated_tables_as_text_from_pdf += table.to_string(index=False)
            #         textual_content_from_pdf_list.append(cumulated_tables_as_text_from_pdf)     

            #         print(f"Processing content from PDF file: {os.path.basename(file.name)}")
#process md files 
            # elif filename.endswith('.md'):
            #     with open(file_path, 'r', encoding='Latin-1') as file:
            #         markdown_content = file.read()
            #         textual_content_from_md_list.append(markdown_content)
                    
            #         # Convert Markdown to plain text
            #         plain_text = markdown.markdown(markdown_content, extensions=['markdown.extensions.tables'])
            #         textual_content_from_md_list.append(plain_text)
                    
            #         print(f"Processing content from Markdown file: {os.path.basename(file.name)}")

#process json files 
            elif filename.endswith('.json'):
                with open(file_path, 'r', encoding='utf-8') as file:
                    json_data = json.load(file)
                    textual_content_from_json_list.extend(entry.get('html', []) for entry in json_data)
                print(f"json ::Name: {os.path.basename(file.name)}")
#process xlsx files 
            elif filename.endswith('.xlsx'):
                wb = openpyxl.load_workbook(file_path)
                for sheet_name in wb.sheetnames:
                    ws = wb[sheet_name]
                    data = []
                    for row in ws.iter_rows(values_only=True):
                        data.append([str(cell) if cell is not None else [] for cell in row])
                    df = pd.DataFrame(data)
                    textual_content_from_xlsx_list.append(str(df))
                print(f"excel ::Name: {os.path.basename(filename)}")
#process word files 
            elif filename.endswith('.docx'):
                doc = Document(file_path)
                for paragraph in doc.paragraphs:
                    textual_content_from_word_list.append(paragraph.text)
                print(f"docx ::Name: {os.path.basename(filename)}")
#process csv files 
            elif filename.endswith('.csv'):
                with open(file_path, 'r', encoding='utf-8') as file:
                    csv_reader = csv.reader(file)
                    for row in csv_reader:
                        textual_content_from_csv_list.append(row)
                print(f"csv ::Name: {os.path.basename(file.name)}")
#process pptx files 
            elif filename.endswith('.pptx'):
                prs = Presentation(file_path)
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, 'text'):
                            textual_content_from_pptx_list.append(shape.text)
                print(f"pptx ::Name: {os.path.basename(filename)}")
        Content=[]
        textual_content_from_all_resources_list.append(' '.join(textual_content_from_txt_list))
        textual_content_from_all_resources_list.append(' '.join(textual_content_from_pdf_list))
        textual_content_from_all_resources_list.append(' '.join(textual_content_from_md_list))
        textual_content_from_all_resources_list.append(' '.join(textual_content_from_json_list))
        textual_content_from_all_resources_list.append(' '.join(textual_content_from_xlsx_list))
        textual_content_from_all_resources_list.append(' '.join(textual_content_from_word_list))
        textual_content_from_all_resources_list.append(' '.join(textual_content_from_csv_list))
        textual_content_from_all_resources_list.append(' '.join(textual_content_from_pptx_list))
        textual_content_from_all_resources_list.append(' '.join(Content))
        text_splitter = RecursiveCharacterTextSplitter(chunk_size=1500, chunk_overlap=250)
        splits = text_splitter.create_documents(textual_content_from_all_resources_list) 
        if splits:
            vectorstore = chroma.Chroma.from_documents(splits, embedding=config.embedding_function, persist_directory="./chroma_db/alldoc")
            retriever = vectorstore.similarity_search(query,k=1)
            global_retriever_results.append(retriever)
        else:
            message = "No content available for processing in "

        return global_retriever_results
    
    
    @tool("Get cources for the user")
    def load_content_of_cource(query: str):
        """Charge et traite le contenu des documents des cources"""
        global_retriever_results=[]
        textual_content_from_all_resources_list = []
        textual_content_from_pdf_list=[]
        textual_content_from_md_list=[]
        textual_content_from_json_list=[]
        textual_content_from_xlsx_list=[]
        textual_content_from_word_list=[]
        textual_content_from_csv_list=[]
        textual_content_from_pptx_list=[]
        textual_content_from_txt_list=[]
        directory_path = config.tuto_path
        # parsed_pdf_files=config.parsed_pdf_files
        # llama_parse_sce.parse_pdf_to_md(directory_path,directory_path,parsed_pdf_files)
        for filename in os.listdir(directory_path):
            file_path = os.path.join(directory_path, filename)
            
#process txt files 
            if filename.endswith('.txt'):
                with open(file_path, 'r', encoding='utf-8') as file:
                    textual_content_from_txt_list.append(file.read())
                    print(f"txt ::Name: {os.path.basename(file.name)}")
#process pdf files 
            # elif filename.endswith('.pdf'):
            #     with open(file_path, 'rb') as file:
            #         pdf_reader = PyPDF2.PdfReader(file)
            #         num_pages = len(pdf_reader.pages)
            #         tables = tabula.read_pdf(file, pages='all', multiple_tables=True)
            #         cumulated_text_from_pdf = ""

            #         for page_num in range(num_pages):
            #             page = pdf_reader.pages[page_num]
            #             cumulated_text_from_pdf += page.extract_text()
            #         textual_content_from_pdf_list.append(cumulated_text_from_pdf)

            #         cumulated_tables_as_text_from_pdf = ""
            #         for table in tables:
            #             if isinstance(table, pd.DataFrame):
            #                 cumulated_tables_as_text_from_pdf += table.to_string(index=False)
            #         textual_content_from_pdf_list.append(cumulated_tables_as_text_from_pdf)     

            #         print(f"Processing content from PDF file: {os.path.basename(file.name)}")
#process md files 
            elif filename.endswith('.md'):
                with open(file_path, 'r', encoding='Latin-1') as file:
                    markdown_content = file.read()
                    textual_content_from_md_list.append(markdown_content)
                    
                    # Convert Markdown to plain text
                    plain_text = markdown.markdown(markdown_content, extensions=['markdown.extensions.tables'])
                    textual_content_from_md_list.append(plain_text)
                    
                    print(f"Processing content from Markdown file: {os.path.basename(file.name)}")

#process json files 
            elif filename.endswith('.json'):
                with open(file_path, 'r', encoding='utf-8') as file:
                    json_data = json.load(file)
                    textual_content_from_json_list.extend(entry.get('html', []) for entry in json_data)
                print(f"json ::Name: {os.path.basename(file.name)}")
#process xlsx files 
            elif filename.endswith('.xlsx'):
                wb = openpyxl.load_workbook(file_path)
                for sheet_name in wb.sheetnames:
                    ws = wb[sheet_name]
                    data = []
                    for row in ws.iter_rows(values_only=True):
                        data.append([str(cell) if cell is not None else [] for cell in row])
                    df = pd.DataFrame(data)
                    textual_content_from_xlsx_list.append(str(df))
                print(f"excel ::Name: {os.path.basename(filename)}")
#process word files 
            elif filename.endswith('.docx'):
                doc = Document(file_path)
                for paragraph in doc.paragraphs:
                    textual_content_from_word_list.append(paragraph.text)
                print(f"docx ::Name: {os.path.basename(filename)}")
#process csv files 
            elif filename.endswith('.csv'):
                with open(file_path, 'r', encoding='utf-8') as file:
                    csv_reader = csv.reader(file)
                    for row in csv_reader:
                        textual_content_from_csv_list.append(row)
                print(f"csv ::Name: {os.path.basename(file.name)}")
#process pptx files 
            elif filename.endswith('.pptx'):
                prs = Presentation(file_path)
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, 'text'):
                            textual_content_from_pptx_list.append(shape.text)
                print(f"pptx ::Name: {os.path.basename(filename)}")
        Content=[]
        textual_content_from_all_resources_list.append(' '.join(textual_content_from_txt_list))
        textual_content_from_all_resources_list.append(' '.join(textual_content_from_pdf_list))
        textual_content_from_all_resources_list.append(' '.join(textual_content_from_md_list))
        textual_content_from_all_resources_list.append(' '.join(textual_content_from_json_list))
        textual_content_from_all_resources_list.append(' '.join(textual_content_from_xlsx_list))
        textual_content_from_all_resources_list.append(' '.join(textual_content_from_word_list))
        textual_content_from_all_resources_list.append(' '.join(textual_content_from_csv_list))
        textual_content_from_all_resources_list.append(' '.join(textual_content_from_pptx_list))
        textual_content_from_all_resources_list.append(' '.join(Content))
        text_splitter = RecursiveCharacterTextSplitter(chunk_size=1500, chunk_overlap=250)
        splits = text_splitter.create_documents(textual_content_from_all_resources_list) 
        if splits:
            vectorstore = chroma.Chroma.from_documents(splits, embedding=config.embedding_function, persist_directory="./chroma_db/alldoc")
            retriever = vectorstore.similarity_search(query,k=10)
            global_retriever_results.append(retriever)
        else:
            message = "No content available for processing in "
        print(global_retriever_results)
        return global_retriever_results
    # print(load_content_of_cource("Abdessamad"))
